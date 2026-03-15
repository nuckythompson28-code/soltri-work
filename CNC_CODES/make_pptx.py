from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ─── 색상 팔레트 ───
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1A, 0x1A, 0x2E)
GRAY      = RGBColor(0x64, 0x74, 0x8B)
RED       = RGBColor(0xEF, 0x44, 0x44)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
YELLOW_BG = RGBColor(0xFF, 0xFA, 0xCD)
RED_BG    = RGBColor(0xFF, 0xEB, 0xEE)
GREEN_BG  = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_BG   = RGBColor(0xE3, 0xF2, 0xFD)
TEAL      = RGBColor(0x0D, 0x94, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, shape_type=MSO_SHAPE.RECTANGLE, rotation=0):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.rotation = rotation
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    try:
        txBox.text_frame.paragraphs[0].space_before = Pt(0)
        txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except:
        pass
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕", line_spacing=1.2):
    """lines: list of (text, font_size, color, bold) or just str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        if isinstance(line, tuple):
            txt, fs, clr, b = line
        else:
            txt, fs, clr, b = line, font_size, color, bold
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(fs)
        run.font.color.rgb = clr
        run.font.bold = b
        run.font.name = font_name
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=NAVY, header_text_color=WHITE, font_size=12):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "맑은 고딕"
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = header_text_color
                else:
                    paragraph.font.color.rgb = BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


def slide_header(slide, title, subtitle=None):
    """Common header bar for content slides"""
    add_shape(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=NAVY)
    add_shape(slide, Inches(0), Inches(1.1), W, Inches(0.06), fill_color=ACCENT)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.6), title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.4), subtitle, font_size=16, color=RGBColor(0xA0,0xBB,0xDD))
    # page accent line
    add_shape(slide, Inches(0), Inches(7.3), W, Inches(0.2), fill_color=NAVY)


# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
# decorative shapes
add_shape(s, Inches(-1), Inches(-1), Inches(6), Inches(6), fill_color=DARK_BLUE, shape_type=MSO_SHAPE.OVAL)
add_shape(s, Inches(9), Inches(3), Inches(6), Inches(6), fill_color=DARK_BLUE, shape_type=MSO_SHAPE.OVAL)
add_shape(s, Inches(0), Inches(3.0), W, Inches(0.06), fill_color=ACCENT)

add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
         "CNC 매크로 프로그램 교육자료", font_size=20, color=ACCENT, bold=True)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(1.0),
         "O0852 프로그램 완전 분석", font_size=44, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.3), Inches(11), Inches(0.5),
         "링형 부품 자동 연속 가공 시스템  |  FANUC 매크로", font_size=20, color=RGBColor(0x94,0xA3,0xB8))

add_multiline(s, Inches(1), Inches(4.5), Inches(5), Inches(2.5), [
    ("프로그램 구조 및 흐름 이해", 18, RGBColor(0xCB,0xD5,0xE1), False),
    ("파라미터 입력 방법 및 변수 해설", 18, RGBColor(0xCB,0xD5,0xE1), False),
    ("가공 사이클 상세 분석", 18, RGBColor(0xCB,0xD5,0xE1), False),
    ("알람 코드 및 안전 주의사항", 18, RGBColor(0xCB,0xD5,0xE1), False),
])

# ═══════════════════════════════════════════
# SLIDE 2: 목차
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "목차", "O0852 교육자료 구성")

items = [
    ("01", "전체 요약", "프로그램이 하는 일 한눈에 보기"),
    ("02", "프로그램 구조", "4개 서브프로그램 관계도"),
    ("03", "입력 파라미터", "작업자가 수정하는 변수 (#101~#123)"),
    ("04", "소재별 이송속도", "CN / RS / CM 자동 설정"),
    ("05", "고급 설정 & 안전체크", "시스템 변수 및 검증 로직"),
    ("06", "가공 사이클 상세", "면삭 → 스텝 → 챔퍼 → 절단"),
    ("07", "오토링크 & 잔재 가공", "자동 소재 인출 시스템"),
    ("08", "알람 코드표", "에러 번호별 원인과 대처"),
    ("09", "기계별 M코드 매핑", "9대 기계 설정표"),
    ("10", "작업자 주의사항", "안전 수칙 및 체크리스트"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.5) + Inches(i * 0.55)
    add_shape(s, Inches(0.8), y, Inches(0.6), Inches(0.42), fill_color=ACCENT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.8), y + Pt(4), Inches(0.6), Inches(0.35), num, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), y + Pt(2), Inches(3.5), Inches(0.4), title, font_size=18, color=BLACK, bold=True)
    add_text(s, Inches(5.2), y + Pt(4), Inches(7), Inches(0.35), desc, font_size=14, color=GRAY)


# ═══════════════════════════════════════════
# SLIDE 3: 전체 요약
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "1. 전체 요약", "O0852는 어떤 프로그램인가?")

add_shape(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2), fill_color=BLUE_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.8), Inches(1.65), Inches(11.5), Inches(0.9),
         "O0852는 CNC 선반에서 파이프 소재를 자동으로 연속 가공하여\n링형 부품(부싱/베어링류)을 대량 생산하는 FANUC 매크로 프로그램입니다.",
         font_size=22, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

# 4 key feature boxes
features = [
    ("매크로 자동화", "변수 입력만으로\n소재/치수/조건 설정", ACCENT),
    ("연속 가공", "면삭→스텝→챔퍼→절단\n사이클 자동 반복", TEAL),
    ("오토링크", "소재 부족 시\n자동 인출 후 재가공", ORANGE),
    ("안전 검증", "30개 이상 알람으로\n입력 오류 사전 차단", RED),
]
for i, (title, desc, color) in enumerate(features):
    x = Inches(0.5 + i * 3.15)
    add_shape(s, x, Inches(3.2), Inches(2.9), Inches(1.8), fill_color=None, line_color=color, line_width=2.5, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(s, x, Inches(3.2), Inches(2.9), Inches(0.55), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x + Inches(0.1), Inches(3.28), Inches(2.7), Inches(0.45), title, font_size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), Inches(3.85), Inches(2.6), Inches(1.0), desc, font_size=14, color=BLACK, align=PP_ALIGN.CENTER)

# current settings summary
add_shape(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.7), fill_color=LIGHT_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.8), Inches(5.5), Inches(3), Inches(0.35), "현재 설정값 기준 가공 요약", font_size=16, color=NAVY, bold=True)

summary_data = [
    ["항목", "소재", "완성품", "생산"],
    ["규격", "CN, OD70 × ID56\n길이 543mm", "OD64.9 × ID58.3\n길이 7.823mm", "1개당 9.843mm 소비\n1회 인출당 약 6개"],
]
add_table(s, Inches(0.8), Inches(5.95), Inches(11.5), Inches(0.95), 2, 4, summary_data,
          col_widths=[Inches(1.5), Inches(3.5), Inches(3.5), Inches(3.0)], font_size=13)


# ═══════════════════════════════════════════
# SLIDE 4: 프로그램 구조
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "2. 프로그램 구조", "4개 서브프로그램의 호출 관계")

# Program boxes
progs = [
    ("O0852", "메인 셋업", "파라미터 입력\n입력값 검증\n기계 설정", NAVY, Inches(0.8)),
    ("O9001", "메인 로직", "길이/개수 계산\n알람 체크\n가공 시작 명령", DARK_BLUE, Inches(3.8)),
    ("O9002", "메인 가공", "면삭 → 스텝절삭\n챔퍼 → 절단\n오토링크 반복", ACCENT, Inches(6.8)),
    ("O9003", "잔재 가공", "남은 소재 인출\n척 길이 재설정\n최종 가공 사이클", TEAL, Inches(9.8)),
]

for prog_name, title, desc, color, x in progs:
    add_shape(s, x, Inches(1.7), Inches(2.5), Inches(3.0), fill_color=None, line_color=color, line_width=2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(s, x, Inches(1.7), Inches(2.5), Inches(0.7), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x, Inches(1.78), Inches(2.5), Inches(0.35), prog_name, font_size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.12), Inches(2.5), Inches(0.3), title, font_size=13, color=RGBColor(0xDD,0xDD,0xFF), align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(2.7), Inches(2.1), Inches(1.8), desc, font_size=14, color=BLACK, align=PP_ALIGN.CENTER)

# arrows
for x in [Inches(3.3), Inches(6.3), Inches(9.3)]:
    add_text(s, x, Inches(2.7), Inches(0.5), Inches(0.5), "→", font_size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# call labels
calls = [("M98 P9001", Inches(3.1)), ("M98 P9002", Inches(6.1)), ("M98 P9003", Inches(9.1))]
for label, x in calls:
    add_text(s, x, Inches(3.2), Inches(0.9), Inches(0.3), label, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Tool summary at bottom
add_shape(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), fill_color=LIGHT_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(1.0), Inches(5.3), Inches(4), Inches(0.35), "사용 공구 정리", font_size=16, color=NAVY, bold=True)

tool_data = [
    ["공구", "용도", "주요 공정", "비고"],
    ["T01", "보링바 (내경 가공)", "황삭, 정삭, 챔퍼", "메인 가공 공구"],
    ["T02", "절단 바이트", "면삭, 절단", "면삭 + 파팅"],
    ["T03", "오토링크용", "소재 인출", "클램프/언클램프"],
]
add_table(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.2), 4, 4, tool_data,
          col_widths=[Inches(1.5), Inches(3.5), Inches(3.3), Inches(3.0)], font_size=12)


# ═══════════════════════════════════════════
# SLIDE 5: 입력 파라미터 - 기본 설정
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "3-1. 입력 파라미터 — 기본 설정", "작업자가 수정하는 변수 #101 ~ #108")

add_shape(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5), fill_color=GREEN_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.7), Inches(1.45), Inches(11.5), Inches(0.4),
         "이 영역만 작업자가 수정합니다 (Line 4~29)", font_size=15, color=GREEN, bold=True)

data = [
    ["변수", "설명", "현재값", "입력 규칙"],
    ["#101", "소재 길이 끝자리", "43", "0 ~ 99 범위"],
    ["#102", "소재 길이 백자리", "500", "0, 100, 200, 300, 400, 500만 가능"],
    ["#103", "척 길이 (물림 길이)", "75", "mm 단위"],
    ["#104", "초기 면삭량", "1", "최대 40까지"],
    ["#105", "소재 종류", "1", "1=CN, 2=RS, 3=CM"],
    ["#106", "가공 유형", "3", "3=챔퍼, 4=챔퍼(확장)"],
    ["#107", "황삭 ON/OFF", "0", "0=OFF, 1=ON"],
    ["#108", "단품 모드", "1", "1=한개씩 가공"],
]
add_table(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(3.8), 9, 4, data,
          col_widths=[Inches(1.5), Inches(3.5), Inches(2.3), Inches(5.0)], font_size=13)

add_shape(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.0), fill_color=YELLOW_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multiline(s, Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.9), [
    ("소재 길이 계산법:  #101 + #102 = 전체 소재 길이", 15, BLACK, True),
    ("예시: #101=43, #102=500  →  500 + 43 = 543mm", 14, GRAY, False),
])


# ═══════════════════════════════════════════
# SLIDE 6: 입력 파라미터 - 치수 데이터
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "3-2. 입력 파라미터 — 치수 데이터", "완성품 규격 변수 #109 ~ #118")

data = [
    ["변수", "설명", "현재값", "단위"],
    ["#109", "소재 외경 (RAW OD)", "70", "mm"],
    ["#110", "소재 내경 (RAW ID)", "56", "mm"],
    ["#111", "완성 외경 (FIN OD)", "64.90", "mm"],
    ["#112", "완성 내경 (FIN ID)", "58.30", "mm"],
    ["#113", "완성 길이 (FIN LENGTH)", "7.823", "mm"],
    ["#114", "절단 바이트 폭 (TIP WIDTH)", "2.02", "mm"],
    ["#115", "내경 R (ID RADIUS)", "0.34", "mm"],
    ["#116", "외경 R (OD RADIUS)", "0.36", "mm"],
    ["#117", "내경 챔퍼 (ID CHAMFER)", "0.37", "mm"],
    ["#118", "외경 챔퍼 (OD CHAMFER)", "0.60", "mm"],
    ["#119", "길이>40mm 확인 플래그", "0", "40초과 시 1 필수"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.2), 12, 4, data,
          col_widths=[Inches(1.2), Inches(3.0), Inches(1.5), Inches(1.8)], font_size=12)

# Visual diagram of the part
add_shape(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.2), fill_color=LIGHT_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(8.5), Inches(1.6), Inches(4.0), Inches(0.4), "부품 단면 개념도", font_size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# Ring shape visualization
add_shape(s, Inches(9.2), Inches(2.5), Inches(2.8), Inches(2.8), fill_color=RGBColor(0xBB,0xDE,0xFB), shape_type=MSO_SHAPE.OVAL)
add_shape(s, Inches(9.7), Inches(3.0), Inches(1.8), Inches(1.8), fill_color=WHITE, shape_type=MSO_SHAPE.OVAL)

add_text(s, Inches(8.5), Inches(5.5), Inches(4.1), Inches(0.3), "← OD: 64.90mm →", font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(5.85), Inches(4.1), Inches(0.3), "← ID: 58.30mm →", font_size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(6.2), Inches(4.1), Inches(0.3), "길이: 7.823mm", font_size=12, color=BLACK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 7: 입력 파라미터 - 절삭 조건
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "3-3. 입력 파라미터 — 절삭 조건", "#120 ~ #123")

data = [
    ["변수", "설명", "현재값", "비고"],
    ["#120", "T01 공구 RPM", "1700", "보링바 회전수"],
    ["#121", "T02 공구 RPM", "1700", "절단 바이트 회전수"],
    ["#122", "인출 거리 (PULL DIST)", "1", "오토링크 시 소재 인출 거리"],
    ["#123", "여유값 (MARGIN)", "0.2", "0~3 범위, 가공 마진"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.3), 5, 4, data,
          col_widths=[Inches(1.5), Inches(3.5), Inches(2.3), Inches(5.0)], font_size=14)

add_shape(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5), fill_color=BLUE_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.7), Inches(4.25), Inches(11.5), Inches(0.4),
         "실제 이송속도(F) = RPM × Feed 계수  (소재별로 자동 계산됩니다)", font_size=15, color=DARK_BLUE, bold=True)


# ═══════════════════════════════════════════
# SLIDE 8: 소재별 이송속도
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "4. 소재별 이송속도 자동 설정", "#105 소재 종류에 따라 이송 계수가 자동 결정됩니다")

data = [
    ["소재", "#124\n(T01 황삭)", "#125\n(T01 스텝)", "#126\n(T01 챔퍼)", "#127\n(T02 절단)", "F값 예시\n(RPM=1700 기준)"],
    ["CN (#105=1)", "0.15", "0.12", "0.15", "0.09", "황삭 F255"],
    ["RS (#105=2)", "0.13", "0.08", "0.05", "0.08", "황삭 F221"],
    ["CM (#105=3)", "0.13", "0.09", "0.09", "0.09", "황삭 F221"],
    ["기타 (DEFAULT)", "0.15", "0.10", "0.10", "0.10", "황삭 F255"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.5), 5, 6, data,
          col_widths=[Inches(2.0), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(3.1)], font_size=13)

add_shape(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), fill_color=LIGHT_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.8), Inches(4.4), Inches(5), Inches(0.35), "이송속도 계산 공식", font_size=17, color=NAVY, bold=True)
add_multiline(s, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0), [
    ("F값 = RPM × Feed 계수", 20, ACCENT, True),
    ("", 8, BLACK, False),
    ("예시 1) T01 CN소재 황삭:  F = 1700 × 0.15 = 255 mm/min", 15, BLACK, False),
    ("예시 2) T01 CN소재 스텝:  F = 1700 × 0.12 = 204 mm/min", 15, BLACK, False),
    ("예시 3) T02 CN소재 절단:  F = 1700 × 0.09 = 153 mm/min", 15, BLACK, False),
    ("", 8, BLACK, False),
    ("RS 소재는 전체적으로 이송이 낮습니다 (경도가 높아 보수적 절삭)", 14, GRAY, False),
])


# ═══════════════════════════════════════════
# SLIDE 9: 고급 설정 & 안전 체크
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "5. 고급 설정 & 안전 체크", "시스템 변수 #500 시리즈 — 작업자 수정 금지 구간")

add_shape(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5), fill_color=RED_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.7), Inches(1.44), Inches(11.5), Inches(0.4),
         "DANGER: STOP EDITING — 이 구간은 시스템 자동 설정입니다. 작업자 수정 금지!", font_size=15, color=RED, bold=True)

data = [
    ["변수", "설명", "현재값", "역할"],
    ["#530", "안전 길이", "12", "척 끝에서 안전거리 확보"],
    ["#531", "척 잔여 최소 길이", "15", "가공 후 최소 15mm 조(jaw) 물림 유지 (소재 이탈 방지)"],
    ["#532", "잔여-안전 보정거리", "1.03", "마지막 오토링크 풀링 시 정확한 소재 위치 계산용 보정값"],
    ["#506", "Z축 황삭 최대 깊이", "55", "한 번에 가공할 최대 Z길이"],
    ["#507", "Z축 황삭 여유 깊이", "1", "황삭 시 Z방향 여유"],
    ["#508", "T01 노즈 R", "0.2", "공구 끝 반경"],
    ["#509", "황삭 간격", "0.4", "X방향 간격"],
    ["#510", "챔퍼 간격", "1", "챔퍼 절삭 시 간격"],
]
add_table(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.0), 9, 4, data,
          col_widths=[Inches(1.5), Inches(3.5), Inches(2.3), Inches(5.0)], font_size=13)

add_shape(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), fill_color=YELLOW_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multiline(s, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.7), [
    ("RS40 / 소경(OD<30mm) 안전 보정", 15, BLACK, True),
    ("RS40 소재 또는 외경 30mm 미만 → #129=15 (잔재보정 15mm 자동 추가, 강성 부족 대비)", 13, GRAY, False),
])


# ═══════════════════════════════════════════
# SLIDE 10: 가공 사이클 흐름도
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "6-1. 가공 사이클 전체 흐름", "O9002 메인 가공 프로세스")

# Flow chart
steps = [
    ("좌표계 설정\nG10 L2 P0 Z[#103]", NAVY),
    ("면삭 (T02)\n끝면 평삭", DARK_BLUE),
    ("스텝절삭 (T01)\n황삭/정삭", ACCENT),
    ("챔퍼 (T01)\nOD-R / ID-R", TEAL),
    ("절단 (T02)\n부품 분리 + M12 카운트", ORANGE),
]

for i, (text, color) in enumerate(steps):
    x = Inches(0.4 + i * 2.58)
    add_shape(s, x, Inches(1.6), Inches(2.3), Inches(1.4), fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, x + Inches(0.1), Inches(1.75), Inches(2.1), Inches(1.1), text, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i in range(4):
    x = Inches(2.7 + i * 2.58)
    add_text(s, x, Inches(1.9), Inches(0.4), Inches(0.5), "→", font_size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# repeat logic
add_shape(s, Inches(0.4), Inches(3.3), Inches(12.5), Inches(0.45), fill_color=GREEN_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.6), Inches(3.33), Inches(12), Inches(0.4),
         "↻ 개수(#517)만큼 반복  →  소재 부족 시 오토링크(N210)  →  면삭부터 재시작",
         font_size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Detailed steps
add_text(s, Inches(0.5), Inches(4.1), Inches(4), Inches(0.35), "주요 계산 변수", font_size=16, color=NAVY, bold=True)

calc_data = [
    ["계산식", "의미", "현재값"],
    ["#505 = #113 + #114", "1개당 가공 단위", "9.843mm"],
    ["#500 = #103 - #530", "유효 가공 길이", "63mm"],
    ["#517 = FIX[#500/#505]", "1회 인출당 가공 개수", "6개"],
    ["#518 = #517 × #505", "1회 총 가공 길이", "59.058mm"],
    ["#140 = #101+#102-#129", "전체 소재 길이(보정)", "543mm"],
]
add_table(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.7), 6, 3, calc_data,
          col_widths=[Inches(4.5), Inches(4.0), Inches(3.8)], font_size=13)


# ═══════════════════════════════════════════
# SLIDE 11: 가공 사이클 상세 — 면삭 & 스텝
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "6-2. 가공 상세 — 면삭 & 스텝절삭", "N100 면삭 → N130 스텝절삭")

# Face cutting section
add_shape(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(2.8), fill_color=None, line_color=DARK_BLUE, line_width=2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.5), fill_color=DARK_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.6), Inches(1.43), Inches(5.5), Inches(0.45), "N100 면삭 (Face Cutting) — T02", font_size=16, color=WHITE, bold=True)
add_multiline(s, Inches(0.6), Inches(2.05), Inches(5.8), Inches(2.0), [
    ("1. G00 W3. T02     ← T02 공구 선택, Z+3 접근", 12, BLACK, False),
    ("2. X[#501]           ← X74(소재외경+4) 위치", 12, BLACK, False),
    ("3. Z0.                ← Z원점(워크 끝면)", 12, BLACK, False),
    ("4. G97 M03 S1700   ← 정속회전, 정회전", 12, BLACK, False),
    ("5. G98 G01 X[#502] ← X52.5까지 면삭", 12, BLACK, False),
    ("   F = 1700 × 0.09 = 153 mm/min", 12, ACCENT, True),
    ("6. G00 X[#111] T01 ← T01로 교환", 12, BLACK, False),
])

# Step cutting section
add_shape(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(2.8), fill_color=None, line_color=ACCENT, line_width=2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.5), fill_color=ACCENT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.1), Inches(1.43), Inches(5.5), Inches(0.45), "N130 스텝절삭 (Step Cutting) — T01", font_size=16, color=WHITE, bold=True)
add_multiline(s, Inches(7.1), Inches(2.05), Inches(5.6), Inches(2.0), [
    ("[황삭] #107=1일 때만 실행", 13, RED, True),
    ("X[#111+0.4] 위치에서 Z방향 직선절삭", 12, BLACK, False),
    ("F = 1700 × 0.15 = 255 mm/min", 12, ACCENT, True),
    ("", 6, BLACK, False),
    ("[정삭] #106=3일 때 실행", 13, TEAL, True),
    ("X[#111] 완성외경에서 Z방향 정밀 절삭", 12, BLACK, False),
    ("F = 1700 × 0.12 = 204 mm/min", 12, ACCENT, True),
    ("M[#134] 보링바 DOWN → G4 X.1 대기 → U0.2 빠짐", 12, BLACK, False),
])

# Chamfer & cutting
add_shape(s, Inches(0.4), Inches(4.5), Inches(6.2), Inches(2.7), fill_color=None, line_color=TEAL, line_width=2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(s, Inches(0.4), Inches(4.5), Inches(6.2), Inches(0.5), fill_color=TEAL, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.6), Inches(4.53), Inches(5.5), Inches(0.45), "N160 챔퍼 가공 (Chamfer) — T01", font_size=16, color=WHITE, bold=True)
add_multiline(s, Inches(0.6), Inches(5.15), Inches(5.8), Inches(1.9), [
    ("TYPE 3 (#106=3): 기본 챔퍼", 13, BLACK, True),
    ("  완성외경 위치 → 챔퍼 시작점 접근", 12, BLACK, False),
    ("  X[#111-(#510+#116)×2] : OD-R 가공", 12, BLACK, False),
    ("  X[#111+(#510+#115)×2] : ID-R 가공", 12, BLACK, False),
    ("", 6, BLACK, False),
    ("TYPE 4 (#106=4): 확장 챔퍼", 13, BLACK, True),
    ("  + 보링바 UP/DOWN 추가 동작", 12, BLACK, False),
    ("  + 완성길이 위치에서 #118/#117 추가 챔퍼", 12, BLACK, False),
])

add_shape(s, Inches(6.9), Inches(4.5), Inches(6.0), Inches(2.7), fill_color=None, line_color=ORANGE, line_width=2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(s, Inches(6.9), Inches(4.5), Inches(6.0), Inches(0.5), fill_color=ORANGE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.1), Inches(4.53), Inches(5.5), Inches(0.45), "N180 절단 (Cutting) — T02", font_size=16, color=WHITE, bold=True)
add_multiline(s, Inches(7.1), Inches(5.15), Inches(5.6), Inches(1.9), [
    ("TYPE 3: 챔퍼 동시 절단", 13, BLACK, True),
    ("  Z[#118-#522] 위치 진입", 12, BLACK, False),
    ("  X[#111] → U-[#118×2] W-[#118] 경사절삭", 12, BLACK, False),
    ("  X[#504] 내경까지 관통 → M12 카운트", 12, BLACK, False),
    ("", 6, BLACK, False),
    ("TYPE 4: 직선 절단", 13, BLACK, True),
    ("  Z-[#522] 위치에서 수직 진입", 12, BLACK, False),
    ("  X[#503] → X[#111] → X[#504] 관통", 12, BLACK, False),
    ("  M12 부품 카운트", 12, RED, True),
])


# ═══════════════════════════════════════════
# SLIDE 12: 오토링크 & 잔재 가공
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "7. 오토링크 & 잔재 가공", "N210 자동 소재 인출 시스템")

# Auto link steps
add_text(s, Inches(0.5), Inches(1.4), Inches(5), Inches(0.35), "오토링크 동작 순서 (N210)", font_size=17, color=NAVY, bold=True)

al_steps = [
    ("1", "M05 스핀들 정지", "안전을 위해 회전 정지"),
    ("2", "G00 T03 → X0", "오토링크 공구 선택, 중심 이동"),
    ("3", "Z[-#522+#122]", "인출 위치로 접근 (F200)"),
    ("4", "M[#132] 오토로더 CLOSE", "오토로더가 소재를 잡음"),
    ("5", "M69 척 UNCLAMP", "척이 소재를 놓음 (2.5초 대기)"),
    ("6", "W[가공길이+여유]", "소재를 앞으로 당김"),
    ("7", "M68 척 CLAMP", "척이 소재를 다시 잡음 (2.5초 대기)"),
    ("8", "M[#131] 오토로더 OPEN", "오토로더 해제"),
    ("9", "G00 Z100. → GOTO 100", "안전위치 복귀 → 면삭부터 재시작"),
]

for i, (num, action, desc) in enumerate(al_steps):
    y = Inches(1.85 + i * 0.52)
    color = RED if "UNCLAMP" in action or "CLAMP" in action else ACCENT
    add_shape(s, Inches(0.5), y, Inches(0.4), Inches(0.4), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_text(s, Inches(0.5), y + Pt(3), Inches(0.4), Inches(0.35), num, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.05), y + Pt(1), Inches(3.5), Inches(0.38), action, font_size=13, color=BLACK, bold=True)
    add_text(s, Inches(4.6), y + Pt(3), Inches(3.0), Inches(0.35), desc, font_size=12, color=GRAY)

# Warning
add_shape(s, Inches(7.8), Inches(1.5), Inches(5.0), Inches(2.5), fill_color=RED_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(8.0), Inches(1.6), Inches(4.6), Inches(0.35), "안전 경고", font_size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
add_multiline(s, Inches(8.0), Inches(2.1), Inches(4.6), Inches(1.8), [
    ("오토링크 동작 중에는", 15, BLACK, True),
    ("척이 열리고 닫히는 구간이 있습니다.", 15, BLACK, True),
    ("", 8, BLACK, False),
    ("이 구간에서 절대로", 15, RED, True),
    ("소재나 공구에 손대지 마세요!", 15, RED, True),
    ("", 8, BLACK, False),
    ("대기 시간(G04)이 충분히 설정되어", 14, GRAY, False),
    ("있으므로 임의로 줄이지 마세요.", 14, GRAY, False),
])

# O9003 section
add_shape(s, Inches(7.8), Inches(4.3), Inches(5.0), Inches(2.8), fill_color=LIGHT_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(8.0), Inches(4.4), Inches(4.6), Inches(0.35), "O9003 잔재 가공", font_size=16, color=NAVY, bold=True)
add_multiline(s, Inches(8.0), Inches(4.9), Inches(4.6), Inches(2.0), [
    ("오토링크 횟수를 모두 사용한 뒤", 14, BLACK, False),
    ("남은 소재를 처리하는 서브프로그램", 14, BLACK, False),
    ("", 8, BLACK, False),
    ("1. 남은 소재 길이 계산 (#540)", 13, BLACK, False),
    ("2. 척 길이 재설정 (#103=#530+#540)", 13, BLACK, False),
    ("3. 좌표계 재설정 (G10 L2 P0)", 13, BLACK, False),
    ("4. 면삭→스텝→챔퍼→절단 반복", 13, BLACK, False),
    ("5. 완료 후 M99 리턴", 13, BLACK, False),
])


# ═══════════════════════════════════════════
# SLIDE 13: 알람 코드표
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "8. 알람 코드표", "에러 발생 시 알람 번호로 원인을 확인하세요")

data = [
    ["알람", "조건", "의미", "조치"],
    ["1", "#140 > 580", "소재 길이 580mm 초과", "#101, #102 확인"],
    ["2", "#109 < #111", "소재 외경 < 완성 외경", "#109, #111 확인"],
    ["3", "#110 > #112", "소재 내경 > 완성 내경", "#110, #112 확인"],
    ["4", "#125 >= 0.14", "T01 스텝 이송 과대", "소재종류(#105) 확인"],
    ["5", "#102 단위 오류", "100단위 아님", "0/100/200/300/400/500"],
    ["6", "#101 > 99", "끝자리 범위 초과", "0~99로 수정"],
    ["7", "#140 <= #103", "소재 길이 부족", "소재 길이 확인"],
    ["10", "#104 > 40", "초기 면삭량 과대", "#104를 40이하로"],
    ["101~", "변수 미입력", "해당 변수가 -9999", "해당 변수 입력"],
    ["203", "OD <= ID", "소재 외경 <= 내경", "#109, #110 확인"],
    ["204", "FIN OD <= ID", "완성 외경 <= 내경", "#111, #112 확인"],
    ["208", "TIP >= FIN LEN", "절단폭 >= 완성길이", "#114, #113 확인"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6), 13, 4, data,
          col_widths=[Inches(1.2), Inches(3.0), Inches(4.0), Inches(4.1)], font_size=12, header_color=RED)


# ═══════════════════════════════════════════
# SLIDE 14: 기계별 M코드 매핑
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "9. 기계별 M코드 매핑", "#130 기계번호에 따른 오토로더/보링바 M코드")

data = [
    ["기계번호", "이름", "AL OPEN\n(#131)", "AL CLOSE\n(#132)", "BR UP\n(#133)", "BR DOWN\n(#134)"],
    ["1", "AL-1", "M56", "M55", "M54", "M53"],
    ["4", "HA-4", "M52", "M51", "M54", "M53"],
    ["5", "AL-5", "M171", "M170", "M53", "M54"],
    ["7", "AL-7", "M64", "M63", "M53", "M54"],
    ["8", "AL-8", "M64", "M63", "M54", "M53"],
    ["9", "AL-9", "M63", "M64", "M53", "M54"],
    ["10", "AL-10", "M64", "M63", "M54", "M53"],
    ["13", "HA-S3", "M56", "M55", "M54", "M53"],
    ["14", "HA-S4", "M56", "M55", "M54", "M53"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5), 10, 6, data,
          col_widths=[Inches(1.5), Inches(2.0), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)], font_size=14)

add_shape(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8), fill_color=BLUE_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_multiline(s, Inches(0.8), Inches(6.25), Inches(11.5), Inches(0.7), [
    ("AL = Auto Loader (자동 소재 공급장치)  |  BR = Boring bar (보링바)", 15, DARK_BLUE, True),
    ("기계 변경 시 #130 값만 해당 번호로 수정하면 M코드가 자동 매핑됩니다", 13, GRAY, False),
])


# ═══════════════════════════════════════════
# SLIDE 15: 작업자 주의사항
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "10. 작업자 주의사항", "안전 수칙 및 필수 체크리스트")

# Must check
add_shape(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(2.8), fill_color=None, line_color=RED, line_width=2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(0.5), fill_color=RED, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.6), Inches(1.43), Inches(5.5), Inches(0.45), "필수 확인 사항", font_size=17, color=WHITE, bold=True)
add_multiline(s, Inches(0.6), Inches(2.05), Inches(5.8), Inches(2.0), [
    ("1. 수정 가능 영역: #101~#123만 수정", 14, BLACK, True),
    ("   → DANGER 이후 구간은 절대 수정 금지", 13, RED, False),
    ("", 6, BLACK, False),
    ("2. #119 확인: 완성길이 > 40mm → #119=1 필수", 14, BLACK, True),
    ("", 6, BLACK, False),
    ("3. #102 입력: 0/100/200/300/400/500만 허용", 14, BLACK, True),
    ("", 6, BLACK, False),
    ("4. #101 입력: 0~99 범위만 허용", 14, BLACK, True),
    ("", 6, BLACK, False),
    ("5. 기계 변경 시 #130만 수정", 14, BLACK, True),
])

# Safety
add_shape(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(2.8), fill_color=None, line_color=ORANGE, line_width=2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.5), fill_color=ORANGE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(7.1), Inches(1.43), Inches(5.5), Inches(0.45), "안전 관련", font_size=17, color=WHITE, bold=True)
add_multiline(s, Inches(7.1), Inches(2.05), Inches(5.6), Inches(2.0), [
    ("1. 오토링크 중 개입 금지", 14, BLACK, True),
    ("   척 개폐 중 소재/공구 접촉 금지!", 13, RED, False),
    ("", 6, BLACK, False),
    ("2. RS40 소재: 안전보정 자동 적용 (15mm)", 14, BLACK, True),
    ("   강성 부족 대비, 임의 해제 금지", 13, GRAY, False),
    ("", 6, BLACK, False),
    ("3. 소경 소재 (OD<30mm): 동일 보정 적용", 14, BLACK, True),
    ("", 6, BLACK, False),
    ("4. 알람 발생 시: 번호 확인 → 코드표 참조", 14, BLACK, True),
    ("   → 해당 변수 수정 후 재시작", 13, GRAY, False),
])

# Checklist
add_shape(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.7), fill_color=GREEN_BG, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(s, Inches(0.6), Inches(4.6), Inches(5), Inches(0.35), "가공 전 체크리스트", font_size=17, color=GREEN, bold=True)

checks = [
    "소재 종류(#105)와 실제 소재 일치 확인",
    "소재 외경/내경(#109,#110)과 실측값 비교",
    "완성 치수(#111~#118)와 도면 대조",
    "기계번호(#130)와 실제 기계 일치 확인",
    "RPM(#120,#121)이 소재/공구에 적합한지 확인",
    "첫 가공 시 단품모드(#108=1)로 시운전 권장",
]
for i, check in enumerate(checks):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = Inches(0.8 + col * 6.2)
    y = Inches(5.05 + row * 0.52)
    add_text(s, x, y, Inches(5.8), Inches(0.4), f"  {check}", font_size=13, color=BLACK)
    add_shape(s, x - Inches(0.15), y + Pt(4), Inches(0.22), Inches(0.22), fill_color=None, line_color=GREEN, line_width=1.5, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)


# ═══════════════════════════════════════════
# SLIDE 16: 마지막 — 요약
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_shape(s, Inches(-1), Inches(-1), Inches(6), Inches(6), fill_color=DARK_BLUE, shape_type=MSO_SHAPE.OVAL)
add_shape(s, Inches(9), Inches(3), Inches(6), Inches(6), fill_color=DARK_BLUE, shape_type=MSO_SHAPE.OVAL)
add_shape(s, Inches(0), Inches(3.0), W, Inches(0.06), fill_color=ACCENT)

add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
         "O0852 프로그램 핵심 정리", font_size=36, color=WHITE, bold=True)

add_multiline(s, Inches(1), Inches(2.1), Inches(11), Inches(0.8), [
    ("파이프 소재 → 링형 부품 자동 연속 가공 매크로 시스템", 20, RGBColor(0x94,0xA3,0xB8), False),
])

summary_items = [
    "작업자는 #101~#123 변수만 수정 (소재, 치수, 조건)",
    "소재 종류(CN/RS/CM)에 따라 이송속도 자동 설정",
    "면삭 → 스텝절삭 → 챔퍼 → 절단 사이클 자동 반복",
    "소재 부족 시 오토링크로 자동 인출 후 재가공",
    "30개 이상 알람으로 입력 오류/안전 문제 사전 차단",
    "9대 기계 M코드 자동 매핑 (#130 변경만으로 전환)",
]
for i, item in enumerate(summary_items):
    y = Inches(3.4 + i * 0.52)
    add_shape(s, Inches(1), y + Pt(5), Inches(0.3), Inches(0.06), fill_color=ACCENT)
    add_text(s, Inches(1.5), y, Inches(10), Inches(0.45), item, font_size=18, color=RGBColor(0xE2,0xE8,0xF0))

add_text(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "질문이 있으시면 언제든 문의하세요", font_size=16, color=RGBColor(0x64,0x74,0x8B), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
output_path = r"C:\Users\admin\Desktop\work\CNC_CODES\O0852_교육자료_v2.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
